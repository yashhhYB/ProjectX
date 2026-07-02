"""
generate_pptx.py — Generate a presentation explaining the Redrob AI Ranking System.
"""

import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation(output_path="Redrob_AI_Approach.pptx"):
    prs = Presentation()

    # Define common layout
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # Slide 1: Title
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Redrob AI Candidate Ranking"
    subtitle.text = "Intelligent Discovery & Matching System\n\nWhat we built, why we built it, and how it works."

    # Slide 2: The Challenge & Philosophy (What & Why)
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "What We Built & Why"
    tf = body_shape.text_frame
    tf.text = "The Challenge:"
    
    p = tf.add_paragraph()
    p.text = "Keyword matching is easily gamed and misses true context."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "The Philosophy (Why we built it this way):"
    p.level = 0
    p.space_before = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "Holistic View: Use all 23 behavioral signals (response rate, notice period, etc.)."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Anti-Gaming: Trust-weighted skills and deep honeypot detection."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Context-Aware: Favor product-company experience and production deployment over consulting-heavy static roles."
    p.level = 1

    # Slide 3: The 9-Component Scoring Engine (How it works)
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "How It Works: 9-Component Scorer"
    tf = body_shape.text_frame
    tf.text = "A balanced, composite evaluation (Sum = 100%):"
    
    components = [
        "25% - Title & Career (Role relevance, seniority, progression)",
        "18% - Skill Match (Proficiency × Duration × Endorsements)",
        "12% - Semantic Similarity (all-mpnet-base-v2 embeddings)",
        "10% - Experience Fit (Gaussian sweet spot around 5-9 yrs)",
        "10% - Behavioral Signals (Response rate, recency, activity)",
        " 8% - Career Depth (Product ratio, shipped-to-prod keywords)",
        " 7% - Assessments (Actual test scores > self-reported)",
        " 5% - Location (India/Pune/Noida focus)",
        " 5% - Education (Relevance & Tier)"
    ]
    
    for comp in components:
        p = tf.add_paragraph()
        p.text = comp
        p.level = 1

    # Slide 4: Career Depth & Behavioral Signals (The Secret Sauce)
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "The Secret Sauce: Deep Analysis"
    tf = body_shape.text_frame
    
    tf.text = "Going beyond standard filtering:"
    
    p = tf.add_paragraph()
    p.text = "Job-Hopper Detection: Average tenure < 18 months incurs a scaling penalty."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Consulting vs. Product: Rewards product company backgrounds as requested by the JD."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Production Keywords: Scans descriptions for 'deployed', 'shipped', 'scale', etc."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Unique Fact-Based Reasoning: The LLM reasoning generator references actual skills, assessment scores, and raises honest concerns instead of using identical templates."
    p.level = 1

    # Slide 5: Enhanced Honeypot Detection
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Security & Validity: 9 Honeypot Checks"
    tf = body_shape.text_frame
    tf.text = "Detecting fakes and bot profiles:"
    
    checks = [
        "Impossible YOE vs Education timeline",
        "Excessive 'Expert' skill inflation for low YOE",
        "Mass zero-duration expert skills",
        "Reversed chronological timelines",
        "Chronological impossibilities (YOE > max possible)",
        "GitHub link integrity anomalies",
        "Admin domain mismatch (Keyword stuffing)",
        "Impossible behavioral bot patterns"
    ]
    
    for check in checks:
        p = tf.add_paragraph()
        p.text = check
        p.level = 1

    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
